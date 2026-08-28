"""
core/views.py
EduDrive platformasining barcha Web view'lari:
auth, dashboard, resource detail/download, upload, profile, telegram link.
"""

import hashlib
import io
import mimetypes
import os
import urllib.parse
import zipfile
from datetime import date

from django.contrib import messages
from django.contrib.auth import authenticate, login, logout
from django.contrib.auth.decorators import login_required
from django.core.paginator import Paginator
from django.db.models import Q, Count, Avg
from django.http import HttpResponseForbidden, JsonResponse, FileResponse, HttpResponse, Http404
from django.shortcuts import render, redirect, get_object_or_404
from django.urls import reverse
from django.utils import timezone
from django.utils.html import escape
from django.views.decorators.clickjacking import xframe_options_exempt
from django.views.decorators.http import require_POST

from .forms import (
    RegisterForm, LoginForm, TelegramLinkForm,
    ResourceUploadForm, ReviewForm, SearchForm, ProfileEditForm,
)
from .models import CustomUser, Category, Resource, Review, Bookmark, Notification


# ---------------------------------------------------------------------------
# AUTHENTICATION
# ---------------------------------------------------------------------------
def register_view(request):
    if request.user.is_authenticated:
        return redirect("home")

    if request.method == "POST":
        form = RegisterForm(request.POST)
        if form.is_valid():
            user = form.save()
            login(request, user)
            user.check_daily_streak()
            messages.success(request, "Xush kelibsiz! Ro'yxatdan muvaffaqiyatli o'tdingiz.")
            return redirect("home")
    else:
        form = RegisterForm()
    return render(request, "register.html", {"form": form})


def login_view(request):
    if request.user.is_authenticated:
        return redirect("home")

    if request.method == "POST":
        form = LoginForm(request.POST)
        if form.is_valid():
            user = authenticate(
                request,
                username=form.cleaned_data["username"],
                password=form.cleaned_data["password"],
            )
            if user is not None:
                login(request, user)
                earned = user.check_daily_streak()
                if earned:
                    messages.success(
                        request,
                        f"Xush kelibsiz, {user.username}! Kunlik faollik uchun +{earned} ball "
                        f"(ketma-ket {user.streak_days} kun) 🔥"
                    )
                else:
                    messages.success(request, f"Xush kelibsiz, {user.username}!")
                next_url = request.GET.get("next") or "home"
                return redirect(next_url)
            messages.error(request, "Login yoki parol noto'g'ri.")
    else:
        form = LoginForm()
    return render(request, "login.html", {"form": form})


def logout_view(request):
    logout(request)
    messages.info(request, "Tizimdan chiqdingiz.")
    return redirect("home")


@login_required
def link_telegram_view(request):
    """
    Foydalanuvchi profilida OTP kod generatsiya qilib, botga yuboradi.
    Bot /link <code> qabul qilib, telegram_id ni shu userga bog'laydi
    (bog'lash logikasi bot/handlers.py da amalga oshiriladi).
    """
    if request.method == "POST":
        form = TelegramLinkForm(request.POST)
        if form.is_valid():
            entered_code = form.cleaned_data["otp_code"]
            if request.user.otp_code == entered_code and request.user.is_otp_valid():
                messages.success(request, "Telegram hisobingiz muvaffaqiyatli ulandi!")
                return redirect("profile")
            messages.error(request, "Kod noto'g'ri yoki muddati o'tgan.")
    else:
        form = TelegramLinkForm()
        # Yangi OTP generatsiya qilamiz, foydalanuvchi buni botga /link buyrug'i bilan yuboradi
        code = TelegramLinkForm.generate_otp(request.user)
    return render(request, "profile.html", {
        "form": form,
        "otp_code": request.user.otp_code,
        "user_obj": request.user,
    })


# ---------------------------------------------------------------------------
# DASHBOARD / HOME
# ---------------------------------------------------------------------------
def home_view(request):
    search_form = SearchForm(request.GET or None)
    resources = Resource.objects.filter(status=Resource.Status.APPROVED).select_related(
        "category", "uploaded_by"
    )

    query = request.GET.get("q", "").strip()
    category_slug = request.GET.get("category", "").strip()
    sort = request.GET.get("sort", "new")

    if query:
        resources = resources.filter(
            Q(title__icontains=query) | Q(description__icontains=query)
        )

    if category_slug:
        resources = resources.filter(category__slug=category_slug)

    if sort == "popular":
        resources = resources.order_by("-download_count")
    elif sort == "rating":
        resources = resources.annotate(avg_rating=Avg("reviews__rating")).order_by("-avg_rating")
    else:
        resources = resources.order_by("-created_at")

    paginator = Paginator(resources, 12)
    page_obj = paginator.get_page(request.GET.get("page"))

    top_downloaded = Resource.objects.filter(
        status=Resource.Status.APPROVED
    ).order_by("-download_count")[:6]

    categories = Category.objects.filter(parent_category__isnull=True).annotate(
        resource_count=Count("resources", filter=Q(resources__status=Resource.Status.APPROVED))
    )

    daily_resource = get_daily_resource()

    context = {
        "page_obj": page_obj,
        "search_form": search_form,
        "top_downloaded": top_downloaded,
        "categories": categories,
        "query": query,
        "active_category": category_slug,
        "active_sort": sort,
        "daily_resource": daily_resource,
    }
    return render(request, "index.html", context)


def get_daily_resource():
    """
    Kun materiali: har kuni bir xil, lekin har kun boshqacha material tanlanadi.
    Sana bo'yicha deterministik tarzda (bugungi sanaga bog'liq hash orqali)
    tasdiqlangan materiallar orasidan tanlab beradi — real random emas, shu
    tufayli sahifa necha marta yangilansa ham kun davomida bir xil ko'rinadi.
    """
    approved_ids = list(
        Resource.objects.filter(status=Resource.Status.APPROVED).values_list("id", flat=True)
    )
    if not approved_ids:
        return None
    seed = int(hashlib.md5(str(date.today()).encode()).hexdigest(), 16)
    chosen_id = approved_ids[seed % len(approved_ids)]
    return Resource.objects.select_related("category", "uploaded_by").get(id=chosen_id)


# ---------------------------------------------------------------------------
# RESOURCE DETAIL & DOWNLOAD
# ---------------------------------------------------------------------------
def resource_detail_view(request, slug):
    resource = get_object_or_404(
        Resource.objects.select_related("category", "uploaded_by"), slug=slug
    )

    # Faqat approved materiallarni hamma ko'ra oladi, boshqasini faqat egasi/admin
    if resource.status != Resource.Status.APPROVED:
        if not request.user.is_authenticated or (
            request.user != resource.uploaded_by and not request.user.is_staff
        ):
            raise Http404("Material topilmadi yoki hali tasdiqlanmagan.")

    resource.increment_view()

    reviews = resource.reviews.select_related("user").all()
    review_form = ReviewForm()
    user_has_bookmarked = False
    user_has_reviewed = False

    if request.user.is_authenticated:
        user_has_bookmarked = Bookmark.objects.filter(
            user=request.user, resource=resource
        ).exists()
        user_has_reviewed = reviews.filter(user=request.user).exists()

        if request.method == "POST" and not user_has_reviewed:
            review_form = ReviewForm(request.POST)
            if review_form.is_valid():
                review = review_form.save(commit=False)
                review.user = request.user
                review.resource = resource
                review.save()
                request.user.add_points(1)
                messages.success(request, "Sharhingiz uchun rahmat! +1 ball qo'shildi.")
                return redirect("resource_detail", slug=slug)

    related = Resource.objects.filter(
        category=resource.category, status=Resource.Status.APPROVED
    ).exclude(id=resource.id)[:4]

    return render(request, "detail.html", {
        "resource": resource,
        "reviews": reviews,
        "review_form": review_form,
        "user_has_bookmarked": user_has_bookmarked,
        "user_has_reviewed": user_has_reviewed,
        "related": related,
    })


def resource_download_view(request, slug):
    """
    Fayl yuklab olinganda download_count +1 qilinadi.
    Avval serverdagi mahalliy fayl (agar bo'lsa) to'g'ridan-to'g'ri beriladi —
    bu eng ishonchli yo'l. Faqat mahalliy fayl yo'q, lekin telegram_file_id
    mavjud bo'lsa (masalan, bot orqali yuklangan material), foydalanuvchi
    Telegram bot orqali faylni olishga yo'naltiriladi.
    """
    resource = get_object_or_404(Resource, slug=slug, status=Resource.Status.APPROVED)
    resource.increment_download()

    if resource.uploaded_by:
        resource.uploaded_by.add_points(2)

    if resource.file:
        return FileResponse(
            resource.file.open("rb"), as_attachment=True,
            filename=resource.file.name.split("/")[-1],
        )

    if resource.telegram_file_id:
        bot_username = os.environ.get("BOT_USERNAME", "EduDriveBot")
        return redirect(f"https://t.me/{bot_username}?start=get_{resource.id}")

    raise Http404("Fayl mavjud emas.")


def _render_office_preview(ext: str, title: str, data: bytes) -> HttpResponse:
    """
    TXT/MD/CSV/DOCX/PPTX fayl baytlaridan demo-ko'rinish (HTML) generatsiya qiladi.
    Ham to'g'ridan-to'g'ri yuklangan resurslar, ham ZIP ichidagi fayllar uchun
    umumiy ishlatiladi.
    """
    if ext in ("txt", "md", "csv"):
        text = data[:4000].decode("utf-8", errors="ignore")
        return HttpResponse(_wrap_simple_page(title, f"<pre>{escape(text)}</pre>"))

    if ext == "docx":
        try:
            import docx
        except ImportError:
            raise Http404("docx ko'rish kutubxonasi o'rnatilmagan.")
        try:
            document = docx.Document(io.BytesIO(data))
        except Exception:
            raise Http404("Fayl ochilmadi — buzilgan yoki himoyalangan bo'lishi mumkin.")

        paragraphs = [p.text for p in document.paragraphs if p.text.strip()][:60]
        body = "".join(f"<p>{escape(p)}</p>" for p in paragraphs) or "<p>Matn topilmadi.</p>"
        return HttpResponse(_wrap_simple_page(title, body))

    if ext == "pptx":
        try:
            from pptx import Presentation
        except ImportError:
            raise Http404("pptx ko'rish kutubxonasi o'rnatilmagan.")
        try:
            presentation = Presentation(io.BytesIO(data))
        except Exception:
            raise Http404("Fayl ochilmadi — buzilgan yoki himoyalangan bo'lishi mumkin.")

        body = ""
        for i, slide in enumerate(presentation.slides, start=1):
            if i > 5:
                break
            texts = [
                shape.text_frame.text for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text.strip()
            ]
            slide_text = "<br>".join(escape(t) for t in texts) or "<i>(matn yo'q)</i>"
            body += f"<h4>Slayd {i}</h4><p>{slide_text}</p>"
        return HttpResponse(_wrap_simple_page(title, body or "<p>Slaydlar topilmadi.</p>"))

    raise Http404("Bu format uchun demo ko'rish mavjud emas.")


@xframe_options_exempt
def resource_preview_view(request, slug):
    """
    Materialni yuklab olmasdan saytda "demo" ko'rish uchun:
    - PDF: avtomatik generatsiya qilingan 2-sahifalik namuna fayl ko'rsatiladi (inline).
    - Rasm: to'liq rasm inline ko'rsatiladi (rasm allaqachon "demo" hisoblanadi).
    - Boshqa formatlar (docx/pptx): demo qo'llab-quvvatlanmaydi, foydalanuvchi
      to'g'ridan-to'g'ri yuklab olishga yo'naltiriladi.
    download_count bu yerda OSHIRILMAYDI — bu faqat ko'rish, yuklab olish emas.
    """
    resource = get_object_or_404(Resource, slug=slug, status=Resource.Status.APPROVED)

    if resource.is_image and resource.file:
        content_type, _ = mimetypes.guess_type(resource.file.name)
        return FileResponse(resource.file.open("rb"), content_type=content_type or "image/jpeg")

    if resource.preview_file:
        return FileResponse(
            resource.preview_file.open("rb"),
            content_type="application/pdf",
            filename=f"demo_{resource.slug}.pdf",
        )

    if resource.is_zip and resource.file:
        try:
            with zipfile.ZipFile(resource.file.open("rb")) as zf:
                entries = [e for e in zf.infolist() if not e.is_dir()]
        except zipfile.BadZipFile:
            raise Http404("Zip fayl buzilgan yoki noto'g'ri formatda.")

        image_exts = ("jpg", "jpeg", "png", "gif", "webp", "bmp")
        viewable_exts = ("pdf", "txt", "md", "csv", "docx", "pptx")
        rows = ""
        for e in entries:
            name = e.filename
            ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
            size_kb = max(e.file_size // 1024, 1)

            if ext in viewable_exts:
                url = reverse("resource_preview_zip_entry", kwargs={
                    "slug": resource.slug, "entry_name": urllib.parse.quote(name, safe=""),
                })
                action = f'<a href="{url}" style="color:#4f46e5;font-weight:600;text-decoration:none;">👁 Ko\'rish</a>'
            elif ext in image_exts:
                action = '<span style="color:#9ca3af;font-size:12px;">Rasmni ko\'rib bo\'lmaydi, iltimos yuklab oling</span>'
            else:
                action = '<span style="color:#9ca3af;font-size:12px;">🔒 Demo yo\'q</span>'

            rows += (
                f"<tr><td>📄 {name}</td>"
                f"<td style='text-align:right;color:#9ca3af'>{size_kb} KB</td>"
                f"<td style='text-align:right'>{action}</td></tr>"
            )

        html = f"""<!DOCTYPE html>
<html lang="uz"><head><meta charset="utf-8">
<style>
  body {{ font-family: -apple-system, sans-serif; margin: 0; padding: 24px; background: #f3f4f6; color: #111827; }}
  h3 {{ color: #4f46e5; margin: 0 0 16px; font-size: 18px; }}
  table {{ width: 100%; border-collapse: collapse; background: #fff; border-radius: 12px; overflow: hidden; }}
  td {{ padding: 10px 14px; border-bottom: 1px solid #e5e7eb; font-size: 14px; }}
  .note {{ color: #9ca3af; font-size: 12px; margin-top: 16px; }}
</style></head>
<body>
  <h3>📦 {resource.title} — arxiv tarkibi ({len(entries)} ta fayl)</h3>
  <table>{rows}</table>
  <p class="note">PDF fayllarni "Ko'rish" tugmasi orqali ochish mumkin. Boshqa formatlar (rasm va h.k.) demo rejimida qo'llab-quvvatlanmaydi — to'liq arxivni yuklab oling.</p>
</body></html>"""
        return HttpResponse(html)

    if resource.is_office_doc and resource.file:
        try:
            data = resource.file.open("rb").read()
        finally:
            resource.file.close()
        return _render_office_preview(resource.file_type.lower(), resource.title, data)

    raise Http404("Bu material uchun demo ko'rish mavjud emas.")


ZIP_ENTRY_PREVIEW_PAGES = 3  # zip ichidagi PDF uchun ko'rsatiladigan sahifalar soni


def _wrap_simple_page(title: str, body_html: str) -> str:
    """ZIP ichidagi matn/docx/pptx demolari uchun umumiy HTML shablon."""
    return f"""<!DOCTYPE html>
<html lang="uz"><head><meta charset="utf-8">
<style>
  body {{ font-family: -apple-system, sans-serif; margin: 0; padding: 24px; background: #f3f4f6; color: #111827; line-height: 1.6; }}
  h3 {{ color: #4f46e5; margin: 0 0 16px; font-size: 16px; word-break: break-all; }}
  .content {{ background: #fff; border-radius: 12px; padding: 20px; font-size: 14px; }}
  pre {{ white-space: pre-wrap; word-break: break-word; font-family: inherit; margin: 0; }}
  h4 {{ color: #4f46e5; margin: 16px 0 6px; }}
  .note {{ color: #9ca3af; font-size: 12px; margin-top: 16px; }}
</style></head>
<body>
  <h3>📄 {escape(title)}</h3>
  <div class="content">{body_html}</div>
  <p class="note">Demo rejimida qisqartirilgan matn ko'rsatilmoqda. To'liq mazmun uchun arxivni yuklab oling.</p>
</body></html>"""


@xframe_options_exempt
def resource_preview_zip_entry_view(request, slug, entry_name):
    """
    ZIP arxiv ichidagi bitta faylni (yuklab olmasdan) ko'rish uchun.
    Qo'llab-quvvatlanadigan formatlar: pdf, txt/md/csv, docx, pptx.
    Rasm va boshqa formatlar uchun bu endpoint chaqirilmaydi.
    """
    resource = get_object_or_404(Resource, slug=slug, status=Resource.Status.APPROVED)

    if not resource.is_zip or not resource.file:
        raise Http404("Material zip fayl emas.")

    entry_name = urllib.parse.unquote(entry_name)
    ext = entry_name.rsplit(".", 1)[-1].lower() if "." in entry_name else ""

    try:
        with zipfile.ZipFile(resource.file.open("rb")) as zf:
            data = zf.read(entry_name)
    except (KeyError, zipfile.BadZipFile):
        raise Http404("Fayl arxiv ichida topilmadi.")

    # --- PDF: birinchi bir necha sahifasini ko'rsatish ---
    if ext == "pdf":
        try:
            from pypdf import PdfReader, PdfWriter
        except ImportError:
            raise Http404("PDF ko'rish kutubxonasi o'rnatilmagan.")

        reader = PdfReader(io.BytesIO(data))
        writer = PdfWriter()
        for i in range(min(ZIP_ENTRY_PREVIEW_PAGES, len(reader.pages))):
            writer.add_page(reader.pages[i])
        buffer = io.BytesIO()
        writer.write(buffer)
        buffer.seek(0)
        return FileResponse(buffer, content_type="application/pdf", filename="demo.pdf")

    # --- TXT / MD / CSV / DOCX / PPTX: umumiy funksiya orqali ---
    if ext in ("txt", "md", "csv", "docx", "pptx"):
        return _render_office_preview(ext, entry_name, data)

    raise Http404("Bu format uchun demo ko'rish mavjud emas.")


@login_required
@require_POST
def toggle_bookmark_view(request, slug):
    resource = get_object_or_404(Resource, slug=slug)
    bookmark, created = Bookmark.objects.get_or_create(user=request.user, resource=resource)
    if not created:
        bookmark.delete()
        return JsonResponse({"bookmarked": False})
    return JsonResponse({"bookmarked": True})


# ---------------------------------------------------------------------------
# UPLOAD
# ---------------------------------------------------------------------------
@login_required
def upload_resource_view(request):
    if request.method == "POST":
        form = ResourceUploadForm(request.POST, request.FILES)
        if form.is_valid():
            resource = form.save(user=request.user)
            notify_staff_new_pending_resource(resource)
            messages.success(
                request,
                "Materialingiz yuborildi! Admin/moderator tasdiqlagach, u sahifada ko'rinadi."
            )
            return redirect("profile")
    else:
        form = ResourceUploadForm()
    return render(request, "upload.html", {"form": form})


def notify_staff_new_pending_resource(resource: Resource):
    """
    Yangi material yuklanganda (status='pending') barcha moderatorlarga
    (is_staff=True) bildirishnoma yuboradi — ular saytdagi qo'ng'iroq
    (bell) orqali ko'rishadi va to'g'ridan-to'g'ri moderatsiya sahifasiga
    o'tib, tasdiqlashi/rad etishi mumkin (/admin ga kirmasdan).
    """
    staff_users = CustomUser.objects.filter(is_staff=True)
    for staff_user in staff_users:
        Notification.objects.create(
            user=staff_user,
            title="🆕 Yangi material moderatsiya kutmoqda",
            message=f"'{resource.title}' materiali ({resource.uploaded_by.username} tomonidan) tasdiqlashni kutmoqda.",
            link=reverse("moderation_dashboard"),
        )


# ---------------------------------------------------------------------------
# PROFILE
# ---------------------------------------------------------------------------
@login_required
def profile_view(request):
    if request.method == "POST":
        edit_form = ProfileEditForm(request.POST, request.FILES, instance=request.user)
        if edit_form.is_valid():
            edit_form.save()
            messages.success(request, "Profil yangilandi!")
            return redirect("profile")
    else:
        edit_form = ProfileEditForm(instance=request.user)

    user_resources = Resource.objects.filter(uploaded_by=request.user).order_by("-created_at")
    bookmarks = Bookmark.objects.filter(user=request.user).select_related("resource")
    notifications = Notification.objects.filter(user=request.user)[:10]

    stats = {
        "total_uploads": user_resources.count(),
        "approved": user_resources.filter(status=Resource.Status.APPROVED).count(),
        "pending": user_resources.filter(status=Resource.Status.PENDING).count(),
        "total_downloads": sum(r.download_count for r in user_resources),
    }

    link_form = TelegramLinkForm()
    otp_code = None
    if not request.user.is_telegram_linked:
        otp_code = TelegramLinkForm.generate_otp(request.user)

    rank = CustomUser.objects.filter(points__gt=request.user.points).count() + 1

    return render(request, "profile.html", {
        "user_obj": request.user,
        "user_resources": user_resources,
        "bookmarks": bookmarks,
        "notifications": notifications,
        "stats": stats,
        "form": link_form,
        "otp_code": otp_code,
        "edit_form": edit_form,
        "rank": rank,
    })


# ---------------------------------------------------------------------------
# LEADERBOARD (Reyting jadvali)
# ---------------------------------------------------------------------------
def leaderboard_view(request):
    top_users = CustomUser.objects.filter(points__gt=0).order_by("-points")[:50]
    return render(request, "leaderboard.html", {"top_users": top_users})


# ---------------------------------------------------------------------------
# LIVE SEARCH SUGGESTIONS (AJAX)
# ---------------------------------------------------------------------------
def search_suggestions_view(request):
    query = request.GET.get("q", "").strip()
    if len(query) < 2:
        return JsonResponse({"results": []})

    resources = Resource.objects.filter(
        Q(title__icontains=query) | Q(description__icontains=query),
        status=Resource.Status.APPROVED,
    )[:6]

    results = [
        {
            "title": r.title,
            "url": reverse("resource_detail", kwargs={"slug": r.slug}),
            "category": r.category.name,
        }
        for r in resources
    ]
    return JsonResponse({"results": results})


# ---------------------------------------------------------------------------
# NOTIFICATIONS (AJAX bell dropdown)
# ---------------------------------------------------------------------------
@login_required
def notifications_dropdown_view(request):
    notifications = Notification.objects.filter(user=request.user).order_by("-created_at")[:8]
    unread_count = Notification.objects.filter(user=request.user, is_read=False).count()
    return JsonResponse({
        "unread_count": unread_count,
        "notifications": [
            {
                "id": n.id, "title": n.title, "message": n.message,
                "link": n.link, "is_read": n.is_read,
                "created_at": timezone.localtime(n.created_at).strftime("%d.%m.%Y %H:%M"),
            }
            for n in notifications
        ],
    })


@login_required
@require_POST
def notifications_mark_read_view(request):
    Notification.objects.filter(user=request.user, is_read=False).update(is_read=True)
    return JsonResponse({"status": "ok"})


@login_required
def moderation_dashboard_view(request):
    """
    /admin ga kirmasdan, to'g'ridan-to'g'ri saytda materiallarni
    tasdiqlash/rad etish uchun panel. Faqat is_staff=True foydalanuvchilar
    (moderator/admin) kira oladi.
    """
    if not request.user.is_staff:
        return HttpResponseForbidden("Bu sahifaga faqat moderatorlar kira oladi.")

    pending_resources = Resource.objects.filter(
        status=Resource.Status.PENDING
    ).select_related("category", "uploaded_by").order_by("created_at")

    return render(request, "moderate.html", {"pending_resources": pending_resources})


@login_required
def moderation_history_view(request):
    """Moderator oxirgi tasdiqlagan/rad etgan materiallar tarixini ko'radi."""
    if not request.user.is_staff:
        return HttpResponseForbidden("Bu sahifaga faqat moderatorlar kira oladi.")

    recent = Resource.objects.exclude(
        status=Resource.Status.PENDING
    ).select_related("category", "uploaded_by").order_by("-approved_at", "-created_at")[:30]

    return render(request, "moderate_history.html", {"recent_resources": recent})


# ---------------------------------------------------------------------------
# CATEGORY
# ---------------------------------------------------------------------------
def category_detail_view(request, slug):
    category = get_object_or_404(Category, slug=slug)
    resources = Resource.objects.filter(
        category=category, status=Resource.Status.APPROVED
    ).select_related("category", "uploaded_by")

    query = request.GET.get("q", "").strip()
    sort = request.GET.get("sort", "new")

    if query:
        resources = resources.filter(
            Q(title__icontains=query) | Q(description__icontains=query)
        )

    if sort == "popular":
        resources = resources.order_by("-download_count")
    elif sort == "rating":
        resources = resources.annotate(avg_rating=Avg("reviews__rating")).order_by("-avg_rating")
    else:
        resources = resources.order_by("-created_at")

    paginator = Paginator(resources, 12)
    page_obj = paginator.get_page(request.GET.get("page"))

    return render(request, "index.html", {
        "page_obj": page_obj,
        "active_category": category.slug,
        "active_sort": sort,
        "query": query,
        "category_obj": category,
        "categories": Category.objects.filter(parent_category__isnull=True).annotate(
            resource_count=Count("resources", filter=Q(resources__status=Resource.Status.APPROVED))
        ),
    })


# ---------------------------------------------------------------------------
# ADMIN MODERATION (Web tomondan tezkor tasdiqlash uchun AJAX endpoint)
# ---------------------------------------------------------------------------
@login_required
@require_POST
def moderate_resource_view(request, slug):
    if not request.user.is_staff:
        return HttpResponseForbidden("Ruxsat yo'q.")

    resource = get_object_or_404(Resource, slug=slug)
    action = request.POST.get("action")

    # Bazadan HAR SAFAR qayta tekshiramiz: agar material allaqachon
    # tasdiqlangan/rad etilgan bo'lsa (masalan, bot orqali allaqachon
    # bajarilgan bo'lsa), qayta ball/bildirishnoma berilmasin.
    if resource.status != Resource.Status.PENDING:
        return JsonResponse({
            "status": resource.status,
            "already_processed": True,
            "message": "Bu material allaqachon ko'rib chiqilgan (boshqa joyda tasdiqlangan/rad etilgan).",
        })

    if action == "approve":
        resource.approve()
        resource.uploaded_by.add_points(5)
        Notification.objects.create(
            user=resource.uploaded_by,
            title="Materialingiz tasdiqlandi ✅",
            message=f"'{resource.title}' materiali endi barchaga ko'rinadi. +5 ball qo'shildi!",
            link=reverse("resource_detail", kwargs={"slug": resource.slug}),
        )
    elif action == "reject":
        reason = request.POST.get("reason", "")
        resource.reject(reason)
        reason_text = reason or "ko'rsatilmagan"
        Notification.objects.create(
            user=resource.uploaded_by,
            title="Materialingiz rad etildi ❌",
            message=f"'{resource.title}' materiali rad etildi. Sabab: {reason_text}",
        )
    else:
        return JsonResponse({"error": "Noto'g'ri amal"}, status=400)

    return JsonResponse({"status": resource.status, "already_processed": False})
